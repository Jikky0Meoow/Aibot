import os
import re
import logging
import random
import time
from datetime import datetime, timedelta
from telegram import (
    Update,
    InputFile,
    Poll,
    InlineKeyboardButton,
    InlineKeyboardMarkup
)
from telegram.ext import (
    Updater,
    CommandHandler,
    MessageHandler,
    Filters,
    CallbackContext,
    CallbackQueryHandler,
    PollAnswerHandler
)
from pdfminer.high_level import extract_text
from pptx import Presentation

# تحميل توكن البوت من متغيرات البيئة
TOKEN = os.getenv("TELEGRAM_TOKEN")
HF_TOKEN = os.getenv("HF_API_TOKEN")

# إعدادات الحدود
MAX_QUESTIONS_PER_FILE = 50
MIN_QUESTIONS = 5
QUESTIONS_PER_BATCH = 5
MAX_FILES_PER_HOUR = 2
MAX_FILES_PER_DAY = 5

# متتبع الاستخدام
user_stats = {}

# إعداد التسجيل
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """استخراج النص من الملفات المدعومة"""
    if file_type == 'pdf':
        return extract_text(file_path)
    elif file_type in ['pptx', 'ppt']:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    return ""

def generate_questions(context: str, num_questions: int) -> list:
    """إنشاء أسئلة باستخدام نموذج Hugging Face"""
    from transformers import pipeline
    qg_pipeline = pipeline(
        "text2text-generation",
        model="mrm8488/t5-base-finetuned-question-generation-ap",
        token=HF_TOKEN
    )
    
    questions = []
    for _ in range(num_questions):
        # تقسيم النص إلى أجزاء لتجنب تجاوز الحد الأقصى
        chunk = " ".join(context.split()[:300])
        result = qg_pipeline(
            f"generate question: {chunk}",
            max_length=128,
            num_return_sequences=1
        )
        if result:
            question = result[0]['generated_text'].strip()
            # إنشاء خيارات متعددة
            options = generate_options(context)
            correct_idx = random.randint(0, 3)
            options[correct_idx] = re.sub(r'\?$', '', question.split('?')[0])[:100]
            questions.append({
                'question': question,
                'options': options,
                'correct_idx': correct_idx
            })
    return questions

def generate_options(context: str) -> list:
    """إنشاء خيارات متعددة من النص"""
    words = [word for word in re.findall(r'\b\w+\b', context) if len(word) > 5]
    if len(words) < 4:
        words += ['Anatomy', 'Physiology', 'Pathology', 'Pharmacology']
    return random.sample(words, 4)

def start(update: Update, context: CallbackContext) -> None:
    """رسالة البدء مع القيود"""
    user_id = update.effective_user.id
    reset_user_stats(user_id)
    
    welcome_msg = (
        "🏥 *مرحباً بك في بوت الأسئلة الطبية!*\n\n"
        "❇️ أرسل ملف PDF أو PPT/PPTX وسأقوم بإنشاء أسئلة طبية بناءً على محتواه\n\n"
        "⚠️ *القيود:*\n"
        f"- {MAX_FILES_PER_HOUR} ملفات في الساعة\n"
        f"- {MAX_FILES_PER_DAY} ملفات في 24 ساعة\n\n"
        "🚫 سيتم حظرك مؤقتاً عند تجاوز الحدود"
    )
    update.message.reply_text(welcome_msg, parse_mode='Markdown')

def handle_document(update: Update, context: CallbackContext) -> None:
    """معالجة الملفات المرسلة"""
    user_id = update.effective_user.id
    reset_user_stats(user_id)
    
    # التحقق من القيود
    if not check_limits(user_id):
        update.message.reply_text(
            "⛔ لقد تجاوزت الحد المسموح:\n"
            f"{MAX_FILES_PER_HOUR} ملفات/ساعة\n"
            f"{MAX_FILES_PER_DAY} ملفات/يوم"
        )
        return
    
    # تحميل الملف
    file = update.message.document
    file_type = file.file_name.split('.')[-1].lower()
    
    if file_type not in ['pdf', 'ppt', 'pptx']:
        update.message.reply_text("❌ نوع الملف غير مدعوم! يرجى إرسال PDF أو PPT/PPTX")
        return
    
    # تحديث الإحصائيات
    user_stats[user_id]['file_count'] += 1
    user_stats[user_id]['last_upload'] = datetime.now()
    
    # معالجة الملف
    progress_msg = update.message.reply_text("📥 جاري تحميل الملف...")
    file_path = f"temp_{user_id}.{file_type}"
    file.get_file().download(file_path)
    
    # استخراج المحتوى
    context.bot.edit_message_text(
        chat_id=update.message.chat_id,
        message_id=progress_msg.message_id,
        text="🔍 جاري معالجة الملف واستخراج المحتوى..."
    )
    
    try:
        text_content = extract_text_from_file(file_path, file_type)
        if not text_content:
            raise ValueError("فشل استخراج النص")
        
        # تقدير عدد الأسئلة
        num_pages = text_content.count('\f') + 1
        max_questions = min(MAX_QUESTIONS_PER_FILE, max(MIN_QUESTIONS, num_pages * 2))
        
        context.bot.edit_message_text(
            chat_id=update.message.chat_id,
            message_id=progress_msg.message_id,
            text=f"✅ تم استخراج {num_pages} صفحة\n"
            f"📊 حدد عدد الأسئلة ({MIN_QUESTIONS}-{max_questions}):"
        )
        
        # حفظ البيانات للمراحل القادمة
        context.user_data['file_text'] = text_content
        context.user_data['max_questions'] = max_questions
        context.user_data['current_questions'] = []
        context.user_data['user_answers'] = []
        context.user_data['score'] = 0
        
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        context.bot.edit_message_text(
            chat_id=update.message.chat_id,
            message_id=progress_msg.message_id,
            text="❌ فشل معالجة الملف. يرجى المحاولة بملف آخر."
        )

def handle_question_count(update: Update, context: CallbackContext) -> None:
    """معالجة عدد الأسئلة المطلوب"""
    try:
        num_questions = int(update.message.text)
        max_questions = context.user_data.get('max_questions', MIN_QUESTIONS)
        
        if not (MIN_QUESTIONS <= num_questions <= max_questions):
            update.message.reply_text(
                f"⚠️ الرجاء إدخال رقم بين {MIN_QUESTIONS} و {max_questions}"
            )
            return
            
        # إنشاء الأسئلة
        text_content = context.user_data['file_text']
        progress_msg = update.message.reply_text("🤖 جاري إنشاء الأسئلة...")
        
        questions = generate_questions(text_content, num_questions)
        context.user_data['total_questions'] = len(questions)
        context.user_data['current_questions'] = questions
        context.user_data['current_index'] = 0
        
        # إرسال الدفعة الأولى
        context.bot.delete_message(
            chat_id=update.message.chat_id,
            message_id=progress_msg.message_id
        )
        send_question_batch(update, context)
        
    except ValueError:
        update.message.reply_text("⚠️ يرجى إدخال رقم صحيح")

def send_question_batch(update: Update, context: CallbackContext) -> None:
    """إرسال مجموعة من الأسئلة"""
    questions = context.user_data['current_questions']
    start_idx = context.user_data['current_index']
    end_idx = start_idx + QUESTIONS_PER_BATCH
    batch = questions[start_idx:end_idx]
    
    for i, q in enumerate(batch):
        poll = Poll(
            question=q['question'][:300],
            options=q['options'],
            type=Poll.QUIZ,
            correct_option_id=q['correct_idx']
        )
        context.bot.send_poll(
            chat_id=update.effective_chat.id,
            question=poll.question,
            options=poll.options,
            type=poll.type,
            correct_option_id=poll.correct_option_id
        )
    
    # التحقق من وجود المزيد
    context.user_data['current_index'] = end_idx
    if end_idx < len(questions):
        keyboard = [[InlineKeyboardButton("👉 المزيد", callback_data='next_batch')]]
        context.bot.send_message(
            chat_id=update.effective_chat.id,
            text=f"🧠 تم إرسال {end_idx} من {len(questions)} سؤال",
            reply_markup=InlineKeyboardMarkup(keyboard)
        )
    else:
        show_results(update, context)

def handle_poll_answer(update: Update, context: CallbackContext) -> None:
    """تسجيل إجابات المستخدم"""
    answer = update.poll_answer
    user_id = answer.user.id
    
    if 'user_answers' not in context.user_data:
        context.user_data['user_answers'] = []
    
    # تخزين الإجابة
    context.user_data['user_answers'].append({
        'question_idx': len(context.user_data['user_answers']),
        'selected': answer.option_ids[0]
    })

def next_batch(update: Update, context: CallbackContext) -> None:
    """إرسال الدفعة التالية من الأسئلة"""
    query = update.callback_query
    query.answer()
    query.delete_message()
    send_question_batch(update, context)

def show_results(update: Update, context: CallbackContext) -> None:
    """عرض النتائج النهائية"""
    total = context.user_data['total_questions']
    correct = sum(1 for ans in context.user_data['user_answers']
                  if ans['selected'] == 
                  context.user_data['current_questions'][ans['question_idx']]['correct_idx'])
    
    score_msg = (
        f"🏁 انتهى الاختبار!\n"
        f"📊 نتيجتك: {correct}/{total}\n"
        f"🎯 الدقة: {round(correct/total*100, 1)}%\n\n"
        "أرسل ملفاً جديداً للبدء مرة أخرى"
    )
    context.bot.send_message(
        chat_id=update.effective_chat.id,
        text=score_msg
    )
    
    # إعادة تعيين البيانات
    context.user_data.clear()

def reset_user_stats(user_id: int) -> None:
    """إعادة تعيين إحصائيات المستخدم"""
    if user_id not in user_stats:
        user_stats[user_id] = {
            'file_count': 0,
            'last_upload': None,
            'daily_count': 0,
            'last_daily_reset': datetime.now()
        }
    
    # إعادة تعيين العد اليومي
    now = datetime.now()
    if now - user_stats[user_id]['last_daily_reset'] > timedelta(hours=24):
        user_stats[user_id]['daily_count'] = 0
        user_stats[user_id]['last_daily_reset'] = now

def check_limits(user_id: int) -> bool:
    """التحقق من حدود الاستخدام"""
    user_data = user_stats.get(user_id)
    if not user_data:
        return True
    
    now = datetime.now()
    
    # التحقق من الحد الساعي
    if user_data['last_upload'] and (now - user_data['last_upload'] < timedelta(hours=1)):
        if user_data['file_count'] >= MAX_FILES_PER_HOUR:
            return False
    
    # التحقق من الحد اليومي
    if user_data['daily_count'] >= MAX_FILES_PER_DAY:
        return False
    
    return True

def main() -> None:
    """تشغيل البوت"""
    updater = Updater(TOKEN)
    dispatcher = updater.dispatcher

    # تسجيل المعالجين
    dispatcher.add_handler(CommandHandler("start", start))
    dispatcher.add_handler(MessageHandler(Filters.document, handle_document))
    dispatcher.add_handler(MessageHandler(Filters.text & ~Filters.command, handle_question_count))
    dispatcher.add_handler(CallbackQueryHandler(next_batch, pattern='^next_batch$'))
    dispatcher.add_handler(PollAnswerHandler(handle_poll_answer))

    # بدء البوت
    updater.start_polling()
    updater.idle()

if __name__ == '__main__':
    main()